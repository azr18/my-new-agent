#!/usr/bin/env python3
"""Generate 'Why Your Business Needs AI Automation Today' pitch deck for mynewagent.ai"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Palette ──────────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x0F, 0x23)   # deep navy
BG_CARD    = RGBColor(0x1A, 0x1A, 0x2E)   # card surface
ACCENT     = RGBColor(0x00, 0xD4, 0xFF)   # electric blue
ACCENT2    = RGBColor(0x00, 0xC9, 0xA7)   # teal
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
WARN       = RGBColor(0xFF, 0x6B, 0x6B)   # coral/red
GREEN      = RGBColor(0x00, 0xE6, 0x76)   # success green
GOLD       = RGBColor(0xFF, 0xD7, 0x00)   # gold highlight

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=18,
                    color=WHITE, bullet_color=ACCENT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u25B8 "  # small right triangle
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = font_name
        # Text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = font_name
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=16, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.name = "Calibri"
        p.font.bold = True
        shape.text_frame.paragraphs[0].space_before = Pt(0)
        shape.text_frame.paragraphs[0].space_after = Pt(0)
        tf.auto_size = None
        # Vertical centering
        from pptx.oxml.ns import qn
        txBody = tf._txBody
        bodyPr = txBody.find(qn('a:bodyPr'))
        bodyPr.set('anchor', 'ctr')
    return shape


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = text


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.6), Inches(4.333), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11.333), Inches(1.2),
                "WHY YOUR BUSINESS NEEDS", font_size=24, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER, bold=False)
    add_textbox(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.2),
                "AI Automation Today", font_size=52, color=WHITE,
                alignment=PP_ALIGN.CENTER, bold=True)
    add_textbox(slide, Inches(1), Inches(3.2), Inches(11.333), Inches(0.8),
                "mynewagent.ai", font_size=28, color=ACCENT,
                alignment=PP_ALIGN.CENTER, bold=False)
    add_textbox(slide, Inches(1), Inches(5.5), Inches(11.333), Inches(0.6),
                "Stop losing money to manual processes. Start building your competitive moat.",
                font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    set_notes(slide, """SPEAKER NOTES — SLIDE 1: TITLE

KEY MESSAGE: Set the tone — this is not a tech lecture, it's a business conversation about money and competitive advantage.

TALKING POINTS:
- Welcome everyone. Thank you for being here.
- "I'm not here to sell you software. I'm here to show you why the businesses that move on AI now will be the ones that dominate their industries in 3 years."
- Quick intro: My New Agent — we're an AI automation consultancy. We build custom automations that replace repetitive manual work.
- "By the end of this presentation, you'll know exactly whether AI automation is right for your business — and I promise, if it's not, I'll tell you that too."

TRANSITION: "Let me start with a question — how many of you feel like AI is something you should be doing something about, but you're not sure what?"
(Pause for hands. This creates engagement and lets you gauge the room.)
""")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 2: AI Adoption Lifecycle
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.7),
                "The AI Adoption Lifecycle", font_size=36, color=WHITE,
                bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11.733), Inches(0.5),
                "Where are we — and what does that mean for your business?",
                font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # Bell curve segments — heights proportional to percentages
    # All bars share bottom edge at 5.6". Heights scaled: 34% = 3.2", others proportional, min 0.8"
    # 2.5% → 0.8", 13.5% → 1.27", 34% → 3.2", 34% → 3.2", 16% → 1.51"
    bottom = 5.6
    segments = [
        ("Innovators\n2.5%",       Inches(1.0),  Inches(bottom - 0.8),  Inches(1.8), Inches(0.8),  RGBColor(0x2D, 0x1B, 0x69)),
        ("Early\nAdopters\n13.5%", Inches(3.0),  Inches(bottom - 1.27), Inches(2.0), Inches(1.27), RGBColor(0x00, 0x50, 0x80)),
        ("Early\nMajority\n34%",   Inches(5.2),  Inches(bottom - 3.2),  Inches(2.2), Inches(3.2),  RGBColor(0x1A, 0x4D, 0x2E)),
        ("Late\nMajority\n34%",    Inches(7.6),  Inches(bottom - 3.2),  Inches(2.2), Inches(3.2),  RGBColor(0x5C, 0x3D, 0x1A)),
        ("Laggards\n16%",          Inches(10.0), Inches(bottom - 1.51), Inches(1.8), Inches(1.51), RGBColor(0x5C, 0x1A, 0x1A)),
    ]
    for label, left, top, w, h, color in segments:
        add_rounded_rect(slide, left, top, w, h, color, label, font_size=14, text_color=WHITE)

    # "YOU ARE HERE" marker
    add_rounded_rect(slide, Inches(2.2), Inches(1.8), Inches(3.5), Inches(0.6),
                     ACCENT, "\u25BC  YOU ARE HERE", font_size=18, text_color=BG_DARK)

    # Two options
    add_textbox(slide, Inches(0.8), Inches(5.9), Inches(5.5), Inches(1.2),
                "\u274C  Option A: Wait & See\nFeel safe today. Get left behind tomorrow.",
                font_size=16, color=WARN, bold=False)
    add_textbox(slide, Inches(6.8), Inches(5.9), Inches(5.5), Inches(1.2),
                "\u2705  Option B: Move Now\n10x–30x efficiency gains. First-mover advantage.",
                font_size=16, color=GREEN, bold=False)

    set_notes(slide, """SPEAKER NOTES — SLIDE 2: THE AI ADOPTION LIFECYCLE

KEY MESSAGE: We are at a once-in-a-generation inflection point. The window to gain first-mover advantage is open RIGHT NOW but it's closing fast.

TALKING POINTS:
- "This is the Technology Adoption Lifecycle — you've probably seen versions of it. It applies to the internet, smartphones, cloud computing — every major tech wave follows this exact pattern."
- "Right now, with AI automation, we are RIGHT HERE" (point to the Early Adopters section). "We've just crossed out of the Innovator phase."
- "That means most businesses haven't moved yet. That's both the risk and the opportunity."

THE TWO OPTIONS:
- "Option A — wait and see. This feels safe. You let someone else figure it out first. The problem? By the time the 'Early Majority' kicks in, the businesses who moved early have already locked in their efficiency gains, their cost advantages, and their customer experience edge. You're now playing catch-up."
- "Option B — move now. Yes, the tech is evolving fast. Yes, you might need to adjust. But the businesses gaining 10x to 30x efficiency today are building a competitive moat that will be almost impossible to overcome."

AUDIENCE ENGAGEMENT:
- "Think about it this way — how many of you wish you'd built a real website in 1998 instead of 2005?"
- Let that land. The analogy is powerful.

TRANSITION: "So the question isn't IF you should adopt AI — it's HOW you do it without blowing your budget or betting on the wrong tech. That's what the next slide is about."
""")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 3: Agility Over Perfection
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.7),
                "Agility Over Perfection", font_size=36, color=WHITE,
                bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11.733), Inches(0.5),
                "The smart way to adopt AI without betting the farm",
                font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # LEFT column: Old Way
    add_rounded_rect(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(0.7),
                     RGBColor(0x5C, 0x1A, 0x1A), "\u274C  THE OLD WAY", font_size=20, text_color=WARN)
    old_way_items = [
        "$50K–$100K on custom-coded software",
        "6–12 month development timeline",
        "Locked into one vendor's proprietary stack",
        "Obsolete when the next AI model drops",
        "Trapped in a sunk-cost fallacy"
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(3.5),
                    old_way_items, font_size=17, color=LIGHT_GRAY, bullet_color=WARN)

    # RIGHT column: New Way
    add_rounded_rect(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.7),
                     RGBColor(0x1A, 0x4D, 0x2E), "\u2705  THE NEW WAY", font_size=20, text_color=GREEN)
    new_way_items = [
        "Lean, agile automation on flexible platforms",
        "Weeks, not months — fast time to value",
        "Built on open tools (n8n, APIs, AI models)",
        "Easy to pivot when better AI arrives",
        "Fraction of the cost, same outcome"
    ]
    add_bullet_list(slide, Inches(7.0), Inches(2.8), Inches(5.5), Inches(3.5),
                    new_way_items, font_size=17, color=LIGHT_GRAY, bullet_color=GREEN)

    # Bottom callout
    add_rounded_rect(slide, Inches(2.5), Inches(6.0), Inches(8.333), Inches(0.8),
                     BG_CARD, "\"The goal isn't perfect software. It's the right outcome at the lowest risk.\"",
                     font_size=17, text_color=ACCENT)

    set_notes(slide, """SPEAKER NOTES — SLIDE 3: AGILITY OVER PERFECTION

KEY MESSAGE: You don't need a $100K software project. You need a lean, flexible automation that delivers the same result at a fraction of the cost — and can evolve with the tech.

TALKING POINTS:
- "Here's the fear I hear from business owners all the time: 'I don't want to drop $50,000 on custom software that's going to be obsolete in a year when the next AI breakthrough happens.' And honestly? That's a SMART fear."
- "The old approach — hiring a dev shop, building proprietary code, 6-12 month timelines — that model is dead for AI. The landscape moves too fast."
- "Our approach is different. We build on agile, third-party platforms. We use tools like n8n — it's an open-source workflow automation platform. Think of it like Zapier on steroids, but you own the infrastructure."

WHY THIS MATTERS:
- "When a better AI model comes out — and it will — we don't throw away your investment. We upgrade the model inside your existing workflow. Your automation gets smarter without starting over."
- "When a new capability becomes available — like an AI that can read invoices or negotiate with vendors — we plug it into your existing system. No rebuild."

OBJECTION HANDLING:
- If someone asks "Why not just use Zapier?" → "Zapier is great for simple tasks. But for business-critical automations that need to handle complex logic, process large volumes, or integrate with custom systems — you need something more robust. That's where we come in."
- If someone asks about data security → "Your data stays on your infrastructure. We don't route your business data through third-party AI services without your explicit approval."

TRANSITION: "So we've talked about WHY now and HOW to do it safely. Let's talk about the number everyone in this room actually cares about — the ROI."
""")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 4: Cost → ROI
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.7),
                "Stop Thinking \"Cost.\" Start Thinking \"ROI.\"", font_size=36, color=WHITE,
                bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11.733), Inches(0.5),
                "Automation isn't a software expense. It's an operational replacement.",
                font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # Big numbers — ROI flow
    add_rounded_rect(slide, Inches(1.0), Inches(2.2), Inches(3.0), Inches(2.0),
                     BG_CARD, "$20,000\nImplementation", font_size=28, text_color=ACCENT)
    # Arrow
    add_textbox(slide, Inches(4.2), Inches(2.8), Inches(1.0), Inches(0.8),
                "\u2794", font_size=48, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, Inches(5.2), Inches(2.2), Inches(3.0), Inches(2.0),
                     BG_CARD, "$60,000\nSaved / Year", font_size=28, text_color=GREEN)
    # Arrow
    add_textbox(slide, Inches(8.4), Inches(2.8), Inches(1.0), Inches(0.8),
                "=", font_size=48, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_rounded_rect(slide, Inches(9.4), Inches(2.2), Inches(3.0), Inches(2.0),
                     RGBColor(0x1A, 0x4D, 0x2E), "$40,000\nNet Year 1", font_size=28, text_color=GOLD)

    # Year-over-year
    add_textbox(slide, Inches(1.0), Inches(4.6), Inches(11.333), Inches(0.5),
                "Every year after Year 1 is pure profit — the automation keeps working, 24/7, with no salary.",
                font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Interactive prompt box
    add_rounded_rect(slide, Inches(1.5), Inches(5.5), Inches(10.333), Inches(1.3),
                     BG_CARD,
                     "\U0001F4AC  QUESTION FOR THE ROOM\n\"What is the most repetitive, time-consuming process "
                     "in your business right now,\nand what does it cost you in employee hours per month?\"",
                     font_size=16, text_color=ACCENT)

    set_notes(slide, """SPEAKER NOTES — SLIDE 4: COST → ROI

KEY MESSAGE: Automation pays for itself — usually in the first year. After that, it's pure profit. Make them do the math on their own business.

TALKING POINTS:
- "Let me reframe how you think about this. Automation is NOT a software expense. It's an operational replacement."
- "Let's say your team spends 40 hours a week on data entry, invoice processing, or lead qualification. That's a full-time salary — let's call it $60,000 a year."
- "We build an automation that handles 80% of that work. Implementation cost: $20,000. In the first year alone, you net $40,000. Every single year after that? That's $60,000 back in your pocket. No sick days. No training. No turnover."

THE INTERACTIVE MOMENT (CRITICAL):
- Read the question on screen out loud.
- "I want you to actually think about this. What's that ONE process that every employee in your company dreads? The one that takes hours, that nobody wants to do, that you KNOW could be done better?"
- Give them 30 seconds of silence. Let them think.
- "Now — what if that process ran itself? What would you do with those freed-up hours? Sell more? Serve clients better? Go home at 5pm for once?"
- This is where you'll see heads nodding. Some will raise hands. Engage with 2-3 responses briefly.

OBJECTION HANDLING:
- "But what about the 20% the automation can't do?" → "Great question. The automation handles the routine 80%. Your people now focus on the 20% that actually requires human judgment — the stuff that grows the business."
- "What if we don't save that much?" → "That's exactly why we start with an audit. We don't guess. We measure your actual hours, actual costs, and we build the business case together before you spend a dime."

TRANSITION: "Now, I know what some of you are thinking — 'This sounds great, but I've been burned before by tech companies promising results.' Let me address that head-on."
""")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 5: Consultative Partner
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.7),
                "A Partner, Not a Vendor", font_size=36, color=WHITE,
                bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11.733), Inches(0.5),
                "Why we say \"no\" more than most agencies say \"yes\"",
                font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # The problem
    add_rounded_rect(slide, Inches(0.8), Inches(1.9), Inches(5.8), Inches(0.6),
                     RGBColor(0x5C, 0x1A, 0x1A), "THE PROBLEM", font_size=18, text_color=WARN)
    problem_items = [
        "Market flooded with \"AI guys\" overpromising",
        "Template solutions sold as custom work",
        "No accountability when it doesn't deliver",
        "You're left holding the bag"
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.7), Inches(5.8), Inches(2.5),
                    problem_items, font_size=17, color=LIGHT_GRAY, bullet_color=WARN)

    # Our approach
    add_rounded_rect(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.6),
                     RGBColor(0x00, 0x50, 0x80), "OUR APPROACH", font_size=18, text_color=ACCENT)
    approach_items = [
        "We audit first — understand before building",
        "If AI isn't right for you, we say so",
        "Custom solutions based on YOUR processes",
        "Ongoing advisory, not just delivery"
    ]
    add_bullet_list(slide, Inches(7.0), Inches(2.7), Inches(5.5), Inches(2.5),
                    approach_items, font_size=17, color=LIGHT_GRAY, bullet_color=ACCENT)

    # Big quote
    add_rounded_rect(slide, Inches(1.5), Inches(5.4), Inches(10.333), Inches(1.4),
                     BG_CARD,
                     "\"If an AI solution isn't right for your current setup,\n"
                     "I will tell you NO. My reputation matters more than a quick sale.\"",
                     font_size=20, text_color=WHITE)

    set_notes(slide, """SPEAKER NOTES — SLIDE 5: CONSULTATIVE PARTNER

KEY MESSAGE: Build trust by being the anti-salesperson. Position yourself as the advisor who prioritizes their business outcome over your revenue.

TALKING POINTS:
- "I want to address something directly. The market right now is FULL of people calling themselves AI experts. Every freelancer with a ChatGPT account is suddenly an 'AI agency.'"
- "And honestly? A lot of them are selling templates. They copy-paste a generic workflow, slap your logo on it, and call it custom. When it breaks — and it WILL break — they're gone."
- "We do it differently. Before we build anything, we do a full process audit. We sit with your team, we map your workflows, we identify where AI actually makes sense — and just as importantly, where it doesn't."

THE POWER LINE:
- Read the quote on screen with conviction: "If an AI solution isn't right for your current setup, I will tell you NO."
- Pause. Let it land.
- "I know that sounds counterintuitive for someone trying to sell you something. But here's the thing — my business grows from referrals. From clients who tell their peers, 'These guys actually delivered.' I'd rather walk away from a bad deal today than damage my reputation."

WHY THIS WORKS:
- In a room of business owners, trust is the #1 currency. They've all been burned by vendors before.
- By positioning yourself as the one who says "no," you become the only one in the room they trust to say "yes."

TRANSITION: "So let me show you exactly how an engagement with us works — step by step."
""")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 6: Our Process
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.7),
                "How We Work With You", font_size=36, color=WHITE,
                bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11.733), Inches(0.5),
                "From first conversation to fully automated operations",
                font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # 5 phase cards
    phases = [
        ("1", "Discover",
         "Learn how your business\nactually runs, day to day.",
         ["Tailored discovery session", "Full process audit", "Priority pain point ranking"],
         RGBColor(0x1E, 0x3A, 0x5F), ACCENT),
        ("2", "Map SOPs",
         "Document every step,\nevery handoff, every bottleneck.",
         ["Step-by-step SOPs", "Time + cost at every step", "Automation ratings per task"],
         RGBColor(0x1E, 0x3A, 0x5F), ACCENT),
        ("3", "Identify\nSavings",
         "Calculate exactly where\ntime and money are leaking.",
         ["Before/after process diagrams", "Dollar value per automation", "Clear ROI breakdown"],
         RGBColor(0x5C, 0x3D, 0x1A), GOLD),
        ("4", "Build It",
         "Build, test, and iterate\nusing your real data.",
         ["Working automations", "Tested on real scenarios", "Error handling built in"],
         RGBColor(0x1A, 0x4D, 0x2E), GREEN),
        ("5", "Go Live",
         "Deploy, train your team,\nand monitor everything.",
         ["Live system, fully deployed", "Team walkthrough + training", "Ongoing monitoring + support"],
         RGBColor(0x1A, 0x4D, 0x2E), GREEN),
    ]

    for i, (num, title, desc, deliverables, bg_color, accent) in enumerate(phases):
        left = Inches(0.3 + i * 2.6)
        card_w = Inches(2.3)

        # Phase card
        add_rounded_rect(slide, left, Inches(1.7), card_w, Inches(0.5),
                         bg_color, f"{num}. {title}", font_size=15, text_color=accent)
        # Description
        add_textbox(slide, left + Inches(0.1), Inches(2.3), Inches(2.1), Inches(0.7),
                    desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

        # "You get:" label
        add_textbox(slide, left + Inches(0.1), Inches(3.1), Inches(2.1), Inches(0.3),
                    "You get:", font_size=12, color=accent, bold=True, alignment=PP_ALIGN.LEFT)

        # Deliverables
        add_bullet_list(slide, left + Inches(0.1), Inches(3.4), Inches(2.1), Inches(2.5),
                        deliverables, font_size=11, color=LIGHT_GRAY, bullet_color=accent)

        # Arrow between phases (except after last)
        if i < 4:
            add_textbox(slide, left + card_w, Inches(1.75), Inches(0.3), Inches(0.4),
                        "\u2794", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom tagline
    add_rounded_rect(slide, Inches(2.0), Inches(6.2), Inches(9.333), Inches(0.7),
                     BG_CARD,
                     "Transparent. Measurable. No surprises.",
                     font_size=17, text_color=WHITE)

    set_notes(slide, """SPEAKER NOTES \u2014 SLIDE 6: OUR PROCESS

KEY MESSAGE: Walk them through the exact engagement model so they know what to expect. No black box, no mystery.

TALKING POINTS:
- "Now let me show you exactly how we work. No mystery, no black box. Five clear phases."

PHASE 1 \u2014 DISCOVER:
- "We start by learning how your business actually runs. Not how you think it runs \u2014 how it ACTUALLY runs day to day."
- "We sit with your team, we watch the workflows, we ask the uncomfortable questions about where time is being wasted."
- "You get a tailored discovery session, a full process audit, and a ranked list of your biggest pain points."

PHASE 2 \u2014 MAP SOPs:
- "Then we document everything. Every step, every handoff, every bottleneck."
- "We create step-by-step SOPs for each process, with the actual time and cost at every step."
- "And we rate every task for automation potential \u2014 some things are easy wins, some aren't worth automating yet."

PHASE 3 \u2014 IDENTIFY SAVINGS:
- "This is where the numbers come in. We calculate exactly where your time and money are leaking."
- "You get before-and-after process diagrams, a dollar value for each automation opportunity, and a clear ROI breakdown."
- "This is where you decide what to build first \u2014 based on hard numbers, not guesses."

PHASE 4 \u2014 BUILD IT:
- "We build, test, and iterate using YOUR real data. Not dummy data, not test scenarios \u2014 your actual business processes."
- "Every automation includes error handling. We plan for the edge cases before they happen."

PHASE 5 \u2014 GO LIVE:
- "We deploy, we train your team, and we monitor everything."
- "This isn't a handoff-and-disappear situation. We stick around to make sure it's working in production."
- "You get ongoing monitoring and support."

LAND THE POINT:
- "The entire process is transparent and measurable. At every phase, you know exactly what you're getting and what it costs. No surprises."

TRANSITION: "And this isn't theoretical. Let me show you real results from businesses we've taken through this exact process."
""")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 7: Real Results (Case Studies)
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.7),
                "Real Results", font_size=36, color=WHITE,
                bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11.733), Inches(0.5),
                "What automation looks like in practice — from businesses like yours",
                font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # Case Study 1: Shipping Automation
    add_rounded_rect(slide, Inches(0.5), Inches(1.8), Inches(3.9), Inches(0.5),
                     RGBColor(0x00, 0x50, 0x80), "PRODUCT COMPANY", font_size=14, text_color=ACCENT)
    add_textbox(slide, Inches(0.5), Inches(2.4), Inches(3.9), Inches(0.7),
                "85%", font_size=48, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(3.0), Inches(3.9), Inches(0.4),
                "reduction in shipping ops time", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
    cs1_items = [
        "5 hrs/day manual \u2192 10 min daily review",
        "Auto label generation + carrier selection",
        "Address validation before printing",
        "~$27,600/year in labor savings"
    ]
    add_bullet_list(slide, Inches(0.7), Inches(3.5), Inches(3.5), Inches(2.5),
                    cs1_items, font_size=14, color=LIGHT_GRAY, bullet_color=ACCENT)

    # Case Study 2: Order Workflow
    add_rounded_rect(slide, Inches(4.7), Inches(1.8), Inches(3.9), Inches(0.5),
                     RGBColor(0x1A, 0x4D, 0x2E), "RETAIL BUSINESS", font_size=14, text_color=GREEN)
    add_textbox(slide, Inches(4.7), Inches(2.4), Inches(3.9), Inches(0.7),
                "ZERO", font_size=48, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.7), Inches(3.0), Inches(3.9), Inches(0.4),
                "data errors in 4+ months", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
    cs2_items = [
        "3-system sync (store, inventory, invoicing)",
        "Auto-validates quantities & totals",
        "Self-corrects when source data changes",
        "~$15\u201319K/year in value"
    ]
    add_bullet_list(slide, Inches(4.9), Inches(3.5), Inches(3.5), Inches(2.5),
                    cs2_items, font_size=14, color=LIGHT_GRAY, bullet_color=GREEN)

    # Case Study 3: Blog Content Pipeline
    add_rounded_rect(slide, Inches(8.9), Inches(1.8), Inches(3.9), Inches(0.5),
                     RGBColor(0x5C, 0x3D, 0x1A), "E-COMMERCE BRAND", font_size=14, text_color=GOLD)
    add_textbox(slide, Inches(8.9), Inches(2.4), Inches(3.9), Inches(0.7),
                "87%", font_size=48, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.9), Inches(3.0), Inches(3.9), Inches(0.4),
                "faster content production", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
    cs3_items = [
        "4\u20136 hrs/post \u2192 30 min review",
        "8-step AI pipeline: plan to publish",
        "SEO-optimized with auto internal linking",
        "52 posts/year, never misses schedule"
    ]
    add_bullet_list(slide, Inches(9.1), Inches(3.5), Inches(3.5), Inches(2.5),
                    cs3_items, font_size=14, color=LIGHT_GRAY, bullet_color=GOLD)

    # Bottom tagline
    add_rounded_rect(slide, Inches(2.0), Inches(6.2), Inches(9.333), Inches(0.7),
                     BG_CARD,
                     "These are single processes for single clients. Most businesses have 3\u20135 like this.",
                     font_size=17, text_color=WHITE)

    set_notes(slide, """SPEAKER NOTES \u2014 SLIDE 7: REAL RESULTS

KEY MESSAGE: These are real results from real businesses. This isn't theoretical \u2014 this is what automation actually delivers.

TALKING POINTS:
- "I've shown you the theory. Now let me show you what this actually looks like. These are three real projects we delivered."

CASE STUDY 1 \u2014 SHIPPING:
- "A product company similar to many of yours was generating shipping labels manually. Copying addresses, picking carriers, downloading labels, pasting tracking numbers back. Five hours a day, every day."
- "We automated the entire flow. Labels generate automatically when orders come in. Carrier is selected based on weight and destination. Addresses are validated BEFORE printing \u2014 so no more wasted labels. Tracking numbers sync back automatically."
- "Result: 85% reduction. Their shipping person now spends 10 minutes reviewing an exception report instead of 5 hours doing data entry. That\u2019s $27,600 a year back in their pocket."

CASE STUDY 2 \u2014 ORDER WORKFLOW:
- "A retail business had data flowing between three systems: their storefront, inventory tracker, and invoicing tool. Every handoff was manual. Someone had to recalculate totals, apply discounts, check quantities. Errors were constant."
- "We built an automated sync with a validation layer. It checks everything before posting. When source data changes, it self-corrects."
- "Result: Zero data errors in over four months. The person who spent 2 hours a day on data entry now spends that time on customer relationships."

CASE STUDY 3 \u2014 CONTENT:
- "An e-commerce brand needed weekly blog posts for SEO. Each post took 4 to 6 hours \u2014 research, writing, images, SEO optimization, formatting."
- "We built an 8-step AI pipeline. Feed it a keyword, and it produces a fully optimized, illustrated, publish-ready blog post."
- "Result: 87% faster. 30 minutes of review instead of 6 hours of work. 52 posts a year, never misses a publish date."

LAND THE POINT:
- Point to the bottom bar: "And here\u2019s the thing \u2014 these are single processes. Most businesses have three to five processes like this. Imagine automating all of them."

TRANSITION: "So now you might be thinking \u2014 'This sounds great, but what happens if it doesn't work? What\u2019s my downside?' Great question. Let me show you how we eliminate your risk entirely."
""")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 8: Total Risk Reversal
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.7),
                "Total Risk Reversal", font_size=36, color=WHITE,
                bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11.733), Inches(0.5),
                "We don't get paid unless we deliver exactly what we promised",
                font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    # Three pillars
    pillars = [
        ("1", "Defined Deliverables",
         "We agree on specific\ntechnical features upfront.\nNo vague promises —\njust hard specs.",
         ACCENT),
        ("2", "Feature Guarantee",
         "If the automation doesn't\nexecute those exact features,\nyou don't pay. Period.\nMoney-back guarantee.",
         GREEN),
        ("3", "Proof of Concept",
         "Before you commit, we show\nyou a working demo tailored\nto YOUR industry. See it\nwork in real-time.",
         GOLD),
    ]

    for i, (num, title, desc, color) in enumerate(pillars):
        left = Inches(0.8 + i * 4.2)
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.3), Inches(2.0), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.color.rgb = BG_DARK
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        from pptx.oxml.ns import qn
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        bodyPr.set('anchor', 'ctr')

        # Title
        add_textbox(slide, left, Inches(2.9), Inches(3.6), Inches(0.5),
                    title, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        # Description
        add_textbox(slide, left, Inches(3.5), Inches(3.6), Inches(2.0),
                    desc, font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom bar
    add_rounded_rect(slide, Inches(2.0), Inches(6.0), Inches(9.333), Inches(0.8),
                     BG_CARD,
                     "Zero risk. Clear deliverables. You only pay for results that work.",
                     font_size=18, text_color=WHITE)

    set_notes(slide, """SPEAKER NOTES — SLIDE 8: TOTAL RISK REVERSAL

KEY MESSAGE: Eliminate the #1 objection — "What if it doesn't work?" — by putting YOUR money where your mouth is.

TALKING POINTS:
- "I know AI can sound like magic. And when something sounds too good to be true, smart business people get skeptical. That's a GOOD instinct."
- "So here's how we eliminate your risk completely."

THREE PILLARS — walk through each:

1. DEFINED DELIVERABLES:
- "Before we write a single line of code, we sit down together and define EXACTLY what the automation will do. Not 'it'll save you time' — specific technical features. 'It will extract invoice data from emails, match it to purchase orders, flag discrepancies, and push clean data to your accounting system.' Black and white."

2. FEATURE GUARANTEE:
- "If the finished automation doesn't do those exact things we agreed on? You don't pay. Full money-back guarantee."
- "This isn't based on your feelings or satisfaction — it's based on hard, measurable technical deliverables. Did it do what we said? Yes or no."
- Let that sink in. This is a powerful differentiator.

3. PROOF OF CONCEPT:
- "Before you even commit to a full engagement, we'll build a simplified working demo that applies to YOUR industry. You'll see the automation running in real-time. Not a slide deck. Not a mockup. A working system."
- "If you don't like what you see? Walk away. No obligation."

AUDIENCE ENGAGEMENT:
- "How many of you have ever paid for software and felt like you didn't get what was promised?" (Hands will go up.)
- "That doesn't happen with us. Because we define success before we start, not after."

TRANSITION: "So — we've covered why now, how to do it safely, the ROI, why you can trust us, the proof, and how we eliminate your risk. Let me wrap up with one question..."
""")

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 9: CTA / Next Steps
    # ════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_DARK)

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.2), Inches(4.333), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    add_textbox(slide, Inches(1), Inches(1.0), Inches(11.333), Inches(1.2),
                "What Will You Do With", font_size=28, color=LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(1.6), Inches(11.333), Inches(1.0),
                "10x More Time?", font_size=52, color=WHITE,
                alignment=PP_ALIGN.CENTER, bold=True)

    add_textbox(slide, Inches(1), Inches(3.0), Inches(11.333), Inches(0.6),
                "Let's find out together.", font_size=24, color=ACCENT,
                alignment=PP_ALIGN.CENTER)

    # CTA boxes
    add_rounded_rect(slide, Inches(1.5), Inches(4.2), Inches(4.8), Inches(1.0),
                     BG_CARD, "\U0001F4E7  ariel@mynewagent.ai", font_size=20, text_color=WHITE)
    add_rounded_rect(slide, Inches(7.0), Inches(4.2), Inches(4.8), Inches(1.0),
                     BG_CARD, "\U0001F310  mynewagent.ai", font_size=20, text_color=WHITE)

    # Free audit offer
    add_rounded_rect(slide, Inches(2.5), Inches(5.8), Inches(8.333), Inches(1.0),
                     ACCENT,
                     "FREE PROCESS AUDIT — Book today and we'll map your #1 automation opportunity",
                     font_size=18, text_color=BG_DARK)

    set_notes(slide, """SPEAKER NOTES — SLIDE 9: CTA / NEXT STEPS

KEY MESSAGE: End with a clear, low-friction call to action. Make it easy for them to take the next step.

TALKING POINTS:
- "Here's the only question that matters: What would you do with 10 times more time?"
- "Would you take on more clients? Would you launch that new product line? Would you finally stop working weekends?"
- Pause. Let them sit with it.

THE OFFER:
- "Here's what I want to offer everyone in this room. A FREE process audit."
- "We'll sit down — 30 minutes, no obligation — and map out your single biggest automation opportunity. We'll show you exactly where you're losing time and money, and what the ROI would look like."
- "If it makes sense to work together, great. If not, you walk away with a clear picture of your own operations. Either way, you win."

CLOSING:
- "My email is on screen. The website is on screen. Come talk to me after this if you have questions."
- "But I'd encourage you — don't wait on this. The businesses that move now are the ones that will be impossible to catch later. Let's make sure that's YOUR business."
- "Thank you."

POST-PRESENTATION:
- Stay at the front of the room. Be approachable.
- Have business cards or a QR code that links directly to a calendar booking page.
- If someone comes up and describes a problem, don't try to solve it on the spot. Instead: "That sounds like a perfect fit for a process audit. Let me get your email and we'll schedule a call this week."
- Follow up within 24 hours with everyone who expressed interest.
""")

    # ── Save ──────────────────────────────────────────────────────────────
    out_dir = os.path.expanduser("~/projects/mynewagent-pitch")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "Why_Your_Business_Needs_AI_Automation.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    return out_path


if __name__ == "__main__":
    build_presentation()
